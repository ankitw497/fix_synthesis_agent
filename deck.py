"""
PowerPoint deck generation for Agent 3 synthesis pipeline.
Implements template binding, single-content enforcement, and exact footers.
"""

import logging
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any, Union
from datetime import datetime
import io
import platform
import subprocess
import shutil

# Stub imports for pptx (would be python-pptx in production)
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import PP_PLACEHOLDER
except ImportError:
    # Stub for development
    Presentation = None
    RGBColor = None
    PP_PLACEHOLDER = None

# Import matplotlib for Figure handling
try:
    from matplotlib.figure import Figure
except ImportError:
    Figure = None
    
from synthesis_agent.config import (
    DECK_FOOTERS, AGENDA_CAVEATS, SynthesisConfig
)
from synthesis_agent.utils import setup_logging

logger = setup_logging("deck")


def _ensure_png_bytes(maybe_fig_or_bytes: Any, dpi: int = 220) -> Optional[bytes]:
    """
    Convert various input types to PNG bytes for embedding in slides.
    
    Args:
        maybe_fig_or_bytes: Figure object, bytes, path string, or None
        dpi: DPI for figure rendering
    
    Returns:
        PNG bytes or None
    """
    if maybe_fig_or_bytes is None:
        return None
    
    if isinstance(maybe_fig_or_bytes, (bytes, bytearray)):
        return bytes(maybe_fig_or_bytes)
    
    if Figure and isinstance(maybe_fig_or_bytes, Figure):
        buf = io.BytesIO()
        maybe_fig_or_bytes.savefig(buf, format="png", dpi=dpi, bbox_inches=None)
        return buf.getvalue()
    
    if isinstance(maybe_fig_or_bytes, str):
        # Assume it's a file path
        try:
            with open(maybe_fig_or_bytes, "rb") as f:
                return f.read()
        except Exception as e:
            logger.warning(f"Could not read file {maybe_fig_or_bytes}: {e}")
            return None
    
    if isinstance(maybe_fig_or_bytes, Path):
        # Handle Path objects
        try:
            return maybe_fig_or_bytes.read_bytes()
        except Exception as e:
            logger.warning(f"Could not read file {maybe_fig_or_bytes}: {e}")
            return None
    
    logger.warning(f"Unsupported figure type: {type(maybe_fig_or_bytes)}")
    return None


def _clear_default_text_artifacts(slide: Any) -> None:
    """Remove vestigial template text boxes that might overlap content."""
    kill_tokens = {"layout", "header", "subheader", "title", "subtitle", "click to add"}
    for shp in list(getattr(slide, "shapes", [])):
        try:
            if getattr(shp, "has_text_frame", False) and hasattr(shp, "text_frame"):
                text = (shp.text_frame.text or "").strip().lower()
                # Remove non-placeholder ghosts by token match
                if (not hasattr(shp, "placeholder_format")) and any(t in text for t in kill_tokens):
                    shp.text_frame.clear()
                # Also clear unused subtitle/content placeholders
                if hasattr(shp, "placeholder_format"):
                    phtype = getattr(shp.placeholder_format, "type", None)
                    if phtype and text in {"subtitle", "content", ""}:
                        shp.text_frame.clear()
        except Exception:
            continue


def _find_content_placeholder(slide: Any) -> Optional[Any]:
    """Return the primary content placeholder on a slide if one exists."""
    title_shape = getattr(slide.shapes, "title", None)
    placeholders = list(getattr(slide, "placeholders", []))
    if PP_PLACEHOLDER is not None:
        preferred = {
            getattr(PP_PLACEHOLDER, "BODY", None),
            getattr(PP_PLACEHOLDER, "CONTENT", None),
        }
        for ph in placeholders:
            try:
                if ph is title_shape:
                    continue
                ph_type = getattr(ph.placeholder_format, "type", None)
                if ph_type in preferred:
                    return ph
            except Exception:
                continue
    for ph in placeholders:
        if ph is title_shape:
            continue
        return ph
    return None


def _remove_unused_content_placeholders(slide: Any) -> None:
    """Delete unused content placeholders so dotted layout boxes disappear."""
    if PP_PLACEHOLDER is None:
        return
    valid = {
        getattr(PP_PLACEHOLDER, "BODY", None),
        getattr(PP_PLACEHOLDER, "CONTENT", None),
    }
    for shp in list(getattr(slide, "shapes", [])):
        try:
            if not getattr(shp, "is_placeholder", False):
                continue
            ph_type = getattr(getattr(shp, "placeholder_format", None), "type", None)
            if ph_type not in valid:
                continue
            if getattr(shp, "has_chart", False) or getattr(shp, "has_table", False):
                continue
            shp._element.getparent().remove(shp._element)
        except Exception:
            continue


def _remove_empty_custom_textboxes(slide: Any) -> None:
    """Remove custom textboxes that have no text content."""
    for shp in list(getattr(slide, "shapes", [])):
        try:
            if getattr(shp, "is_placeholder", False):
                continue
            if not getattr(shp, "has_text_frame", False):
                continue
            text = (shp.text_frame.text or "").strip()
            if not text:
                shp._element.getparent().remove(shp._element)
        except Exception:
            continue



from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

EMU_PER_INCH = 914400


def _apply_dynamic_number_format(chart, spec):
    """Apply a compact number format based on data magnitude."""
    if spec.get("y_is_percent"):
        chart.value_axis.tick_labels.number_format = "0%"
        return
    ymax = 0
    for s in spec.get("series", []):
        for v in s.get("values", []) or []:
            if v is not None:
                try:
                    ymax = max(ymax, abs(float(v)))
                except Exception:
                    continue
    if ymax >= 1_000_000_000:
        fmt = '#,##0,,,"B"'
    elif ymax >= 1_000_000:
        fmt = '#,##0,,"M"'
    elif ymax >= 1_000:
        fmt = '#,##0,"K"'
    else:
        fmt = '#,##0'
    chart.value_axis.tick_labels.number_format = fmt


def _render_chart_spec(slide, placeholder_or_slide, spec, config):
    cd = ChartData(); cd.categories = spec["categories"]
    for s in spec["series"]:
        cd.add_series(s["name"], s["values"])
    kind = spec["kind"]
    ctype_map = {
        "line": "LINE",  # marker-less line charts for trends
        "grouped_bar": "COLUMN_CLUSTERED",
        "stacked_col": "COLUMN_STACKED",
        "combo": "COLUMN_CLUSTERED",
    }
    ctype = getattr(XL_CHART_TYPE, ctype_map.get(kind, "COLUMN_CLUSTERED"))
    if hasattr(placeholder_or_slide, "insert_chart"):
        chart_shape, chart = placeholder_or_slide.insert_chart(ctype, cd)
    else:
        from pptx.util import Inches
        left, top, width, height = spec.get("box", (Inches(1), Inches(1.6), Inches(8), Inches(3.6)))
        chart_shape = slide.shapes.add_chart(ctype, left, top, width, height, cd)
        chart = chart_shape.chart

    # Chart title (optional)
    if spec.get("title") and getattr(config.features, "chart_title_inside", False):
        chart.has_title = True
        chart.chart_title.text_frame.text = spec["title"]

    # Axis titles and formatting
    if getattr(chart, "category_axis", None):
        if spec.get("x_axis_title"):
            chart.category_axis.has_title = True
            chart.category_axis.axis_title.text_frame.text = spec["x_axis_title"]
        try:
            chart.category_axis.tick_labels.font.size = Pt(getattr(config.render, "category_font_pt", 9))
        except Exception:
            pass
        try:
            dense_threshold = getattr(config.render, "category_dense_threshold", 12)
            if len(spec.get("categories", [])) > dense_threshold:
                chart.category_axis.tick_label_spacing = 2
        except Exception:
            pass

    if getattr(chart, "value_axis", None):
        _apply_dynamic_number_format(chart, spec)
        if spec.get("y_axis_title"):
            chart.value_axis.has_title = True
            chart.value_axis.axis_title.text_frame.text = spec["y_axis_title"]
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.has_minor_gridlines = False
        try:
            chart.value_axis.tick_labels.font.size = Pt(getattr(config.render, "value_font_pt", 9))
        except Exception:
            pass
        try:
            from pptx.dml.color import RGBColor
            line = chart.value_axis.major_gridlines.format.line
            line.solid()
            hx = getattr(config.render, "gridline_color", "#E0E0E0")
            r, g, b = int(hx[1:3], 16), int(hx[3:5], 16), int(hx[5:7], 16)
            line.fore_color.rgb = RGBColor(r, g, b)
            line.width = Pt(0.2)
        except Exception:
            pass

    co = dict(spec.get("color_overrides", {}) or {})
    try:
        from synthesis_agent.charts import _norm_tier_key  # Local import to avoid cycles
    except Exception:
        _norm_tier_key = lambda value: str(value or "").strip().upper().replace("-", " ").replace("/", " ").replace(" ", "_")
    try:
        from synthesis_agent.config import STYLE_COLORS
    except Exception:
        STYLE_COLORS = {}
    for series in spec.get("series", []) or []:
        name = series.get("name")
        if not name:
            continue
        key = _norm_tier_key(name)
        if key in STYLE_COLORS and name not in co:
            co[name] = STYLE_COLORS[key]
    for idx, s in enumerate(chart.series):
        name = spec["series"][idx]["name"]
        hx = co.get(name)
        if hx:
            from pptx.dml.color import RGBColor
            r, g, b = int(hx[1:3],16), int(hx[3:5],16), int(hx[5:7],16)
            s.format.fill.solid(); s.format.fill.fore_color.rgb = RGBColor(r,g,b)
    # Hero-first styling and legends
    try:
        chart.chart_area.format.fill.background()
        chart.plot_area.format.fill.background()
        chart.chart_area.format.line.fill.background()
        chart.plot_area.format.line.fill.background()
    except Exception:
        pass

    series_count = len(spec.get("series", []))
    kind = spec.get("kind")
    show_legend = False
    if kind == "line" and series_count >= 1:
        show_legend = True
    elif kind in ("grouped_bar", "stacked_col") and series_count > 1:
        show_legend = True
    if show_legend:
        chart.has_legend = True
        position_name = str(getattr(config.render, "legend_position", "BOTTOM") or "BOTTOM").upper()
        legend_position = getattr(XL_LEGEND_POSITION, position_name, XL_LEGEND_POSITION.BOTTOM)
        chart.legend.position = legend_position
        chart.legend.include_in_layout = True
        try:
            chart.legend.font.size = Pt(getattr(config.render, "legend_font_pt", 11))
        except Exception:
            pass

    try:
        plot = chart.plots[0]
        if spec.get("kind") in ("grouped_bar", "stacked_col"):
            plot.gap_width = 90
            plot.overlap = 0
    except Exception:
        pass

    if kind == "combo":
        # python-pptx 1.0.x lacks public API for secondary axis assignment.
        # Combo specs are handled separately via two-chart fallback or
        # template-based combos. Here we simply skip secondary-axis logic.
        pass
    return chart_shape, chart


def add_chart_slide(prs, title, subtitle, spec, config):
    """Render a single chart spec on a slide."""
    from pptx.util import Inches
    layout = fuzzy_map_layout(prs, desired_layout="Title and Content", config=config)
    slide = prs.slides.add_slide(layout)
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape:
        title_shape.text = title or ""
        try:
            tf = title_shape.text_frame
            if tf.paragraphs:
                tf.paragraphs[0].font.size = Pt(20)
        except Exception:
            pass

    slide_width_in = prs.slide_width / EMU_PER_INCH if getattr(prs, "slide_width", None) else 13.333
    slide_height_in = prs.slide_height / EMU_PER_INCH if getattr(prs, "slide_height", None) else 7.5

    left_in = getattr(config.render, "single_chart_left_in", 0.35)
    right_in = getattr(config.render, "single_chart_right_in", left_in)
    bottom_margin_in = getattr(config.render, "single_chart_bottom_margin_in", 0.6)
    subtitle_gap_in = getattr(config.render, "subtitle_gap_in", 0.10)
    subtitle_height_in = getattr(config.render, "subtitle_height_in", 0.45)
    chart_gap_in = getattr(config.render, "chart_gap_in", 0.05)

    if title_shape:
        title_bottom_in = (title_shape.top + title_shape.height) / EMU_PER_INCH
    else:
        title_bottom_in = getattr(config.render, "title_bottom_in", 0.6)

    chart_width_in = slide_width_in - left_in - right_in
    min_width = getattr(config.render, "single_chart_min_width_in", 6.0)
    desired_width = getattr(config.render, "single_chart_width_in", None)
    if desired_width:
        chart_width_in = min(chart_width_in, desired_width)
    if chart_width_in < min_width:
        chart_width_in = min_width

    subtitle_top_in = title_bottom_in + subtitle_gap_in
    subtitle_bottom_in = subtitle_top_in + subtitle_height_in
    chart_top_in = subtitle_bottom_in + chart_gap_in

    available_height = slide_height_in - chart_top_in - bottom_margin_in
    default_height = getattr(config.render, "single_chart_height_in", 5.9)
    min_height = getattr(config.render, "single_chart_min_height_in", 3.5)
    fill_height = getattr(config.render, "single_chart_fill_height", True)
    if available_height > 0:
        if fill_height:
            chart_height_in = max(min_height, available_height)
        else:
            chart_height_in = max(min_height, min(default_height, available_height))
    else:
        chart_height_in = max(min_height, default_height)

    placeholder = _find_content_placeholder(slide)
    target = placeholder if placeholder is not None and hasattr(placeholder, "insert_chart") else slide
    chart_spec = dict(spec)
    chart_spec["box"] = (
        Inches(left_in),
        Inches(chart_top_in),
        Inches(chart_width_in),
        Inches(chart_height_in),
    )
    if getattr(config.features, "chart_title_inside", False):
        if not chart_spec.get("title"):
            chart_spec["title"] = subtitle or spec.get("title")
    else:
        chart_spec["title"] = None
    chart_shape, _ = _render_chart_spec(slide, target, chart_spec, config)
    try:
        if chart_shape is not None:
            chart_shape.left = Inches(left_in)
            chart_shape.top = Inches(chart_top_in)
            chart_shape.width = Inches(chart_width_in)
            chart_shape.height = Inches(chart_height_in)
    except Exception:
        pass

    if getattr(config.features, "show_subtitle", False) and subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(left_in),
            Inches(subtitle_top_in),
            Inches(chart_width_in),
            Inches(subtitle_height_in),
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        try:
            if subtitle_frame.paragraphs:
                subtitle_frame.paragraphs[0].font.size = Pt(12)
        except Exception:
            pass
    _clear_default_text_artifacts(slide)
    _remove_unused_content_placeholders(slide)
    _remove_empty_custom_textboxes(slide)
    return slide


def add_two_chart_slide(prs, title, subtitle, top_spec, bottom_spec, config):
    from pptx.util import Inches

    layout = fuzzy_map_layout(prs, desired_layout="Title and Content", config=config)
    slide = prs.slides.add_slide(layout)

    title_shape = getattr(slide.shapes, "title", None)
    if title_shape:
        title_shape.text = title or ""
        try:
            tf = title_shape.text_frame
            if tf.paragraphs:
                tf.paragraphs[0].font.size = Pt(20)
        except Exception:
            pass

    slide_width_in = prs.slide_width / EMU_PER_INCH if getattr(prs, "slide_width", None) else 13.333
    slide_height_in = prs.slide_height / EMU_PER_INCH if getattr(prs, "slide_height", None) else 7.5

    left_in = getattr(config.render, "single_chart_left_in", 0.35)
    right_in = getattr(config.render, "single_chart_right_in", left_in)
    bottom_margin_in = getattr(config.render, "single_chart_bottom_margin_in", 0.6)
    subtitle_gap_in = getattr(config.render, "subtitle_gap_in", 0.10)
    subtitle_height_in = getattr(config.render, "subtitle_height_in", 0.45)
    chart_gap_in = getattr(config.render, "chart_gap_in", 0.05)
    inter_chart_gap_in = getattr(config.render, "two_chart_gap_in", 0.25)
    min_height = getattr(config.render, "single_chart_min_height_in", 3.5)

    chart_width_in = slide_width_in - left_in - right_in
    min_width = getattr(config.render, "single_chart_min_width_in", 6.0)
    desired_width = getattr(config.render, "single_chart_width_in", None)
    if desired_width:
        chart_width_in = min(chart_width_in, desired_width)
    if chart_width_in < min_width:
        chart_width_in = min_width

    if title_shape:
        title_bottom_in = (title_shape.top + title_shape.height) / EMU_PER_INCH
    else:
        title_bottom_in = getattr(config.render, "title_bottom_in", 0.6)

    subtitle_top_in = title_bottom_in + subtitle_gap_in
    subtitle_bottom_in = subtitle_top_in + subtitle_height_in
    first_chart_top_in = subtitle_bottom_in + chart_gap_in

    available_height = slide_height_in - first_chart_top_in - bottom_margin_in
    # Divide the remaining vertical space between the two charts with a gap
    chart_height_in = max(min_height, (available_height - inter_chart_gap_in) / 2) if available_height > 0 else min_height

    if getattr(config.features, "show_subtitle", False) and subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(left_in),
            Inches(subtitle_top_in),
            Inches(chart_width_in),
            Inches(subtitle_height_in),
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        try:
            if subtitle_frame.paragraphs:
                subtitle_frame.paragraphs[0].font.size = Pt(12)
        except Exception:
            pass

    # Render top chart
    top_chart_spec = dict(top_spec)
    top_chart_spec["box"] = (
        Inches(left_in),
        Inches(first_chart_top_in),
        Inches(chart_width_in),
        Inches(chart_height_in),
    )
    if getattr(config.features, "chart_title_inside", False):
        if not top_chart_spec.get("title"):
            top_chart_spec["title"] = top_spec.get("title") or subtitle
    else:
        top_chart_spec["title"] = None
    _render_chart_spec(slide, slide, top_chart_spec, config)

    # Render bottom chart
    second_chart_top_in = first_chart_top_in + chart_height_in + inter_chart_gap_in
    bottom_chart_spec = dict(bottom_spec)
    bottom_chart_spec["box"] = (
        Inches(left_in),
        Inches(second_chart_top_in),
        Inches(chart_width_in),
        Inches(chart_height_in),
    )
    if getattr(config.features, "chart_title_inside", False):
        if not bottom_chart_spec.get("title"):
            bottom_chart_spec["title"] = bottom_spec.get("title") or subtitle
    else:
        bottom_chart_spec["title"] = None
    _render_chart_spec(slide, slide, bottom_chart_spec, config)

    _clear_default_text_artifacts(slide)
    _remove_unused_content_placeholders(slide)
    _remove_empty_custom_textboxes(slide)
    return slide


def add_combo_slide(prs, title, subtitle, combo_spec, config):
    """Render a combo spec using configured strategy."""
    if getattr(config.features, "combo_strategy", "two_charts") == "two_charts":
        primary = [s for s in combo_spec.get("series", []) if not s.get("y2")]
        secondary = [s for s in combo_spec.get("series", []) if s.get("y2")]
        top_spec = {**combo_spec, "kind": combo_spec.get("primary_kind", "grouped_bar"), "series": primary}
        bottom_spec = {**combo_spec, "kind": "line", "series": secondary, "title": None}
        return add_two_chart_slide(prs, title, subtitle, top_spec, bottom_spec, config)
    else:
        from pptx.util import Inches
        layout = fuzzy_map_layout(prs, desired_layout="Title and Content", config=config)
        slide = prs.slides.add_slide(layout)
        if getattr(slide.shapes, "title", None):
            slide.shapes.title.text = title or ""
        placeholder = slide.shapes.placeholders[1] if len(slide.shapes.placeholders) > 1 else slide
        _render_chart_spec(slide, placeholder, combo_spec, config)
        tx = slide.shapes.add_textbox(Inches(1), Inches(0.9), Inches(8), Inches(0.4))
        tx.text_frame.text = subtitle or ""
        _clear_default_text_artifacts(slide)
        return slide

def scan_templates_dir(template_dir: Union[str, Path]) -> List[Path]:
    """
    Scan directory for .pptx templates.
    
    Args:
        template_dir: Directory to scan
    
    Returns:
        List of template file paths
    """
    template_path = Path(template_dir)
    
    if not template_path.exists():
        logger.warning(f"Template directory does not exist: {template_dir}")
        return []
    
    templates = list(template_path.glob("*.pptx"))
    logger.info(f"Found {len(templates)} templates in {template_dir}")
    
    return templates


def bind_template(preferred_name: Optional[str],
                 template_dir: Union[str, Path],
                 config: SynthesisConfig,
                 fallback_name: Optional[str] = None) -> Tuple[Any, Dict[str, Any]]:
    """
    Bind to template with preferred→fallback→newest logic.
    Log template readiness report.
    
    Args:
        preferred_name: Preferred template name
        template_dir: Template directory
        config: Synthesis configuration
    
    Returns:
        Tuple of (Presentation, readiness_report)
    """
    readiness_report = {
        'template_used': None,
        'layouts_available': [],
        'single_content_layouts': [],
        'dual_content_layouts': [],
        'missing_layouts': [],
        'warnings': []
    }
    
    # Find templates
    templates = scan_templates_dir(template_dir)
    
    if not templates:
        logger.warning("No templates found, using minimal built-in layout")
        readiness_report['warnings'].append("No templates found, using built-in fallback")
        readiness_report['template_used'] = "built-in"
        return _create_minimal_presentation(), readiness_report
    
    # Select template
    selected_template = None

    if preferred_name:
        for template in templates:
            if preferred_name.lower() in template.name.lower():
                selected_template = template
                break

    if not selected_template and fallback_name:
        for template in templates:
            if fallback_name.lower() in template.name.lower():
                selected_template = template
                readiness_report['warnings'].append(
                    f"Preferred template '{preferred_name}' not found, using fallback '{template.name}'")
                break

    if not selected_template:
        templates.sort(key=lambda x: x.stat().st_mtime, reverse=True)
        selected_template = templates[0]
        readiness_report['warnings'].append(
            f"Preferred template '{preferred_name}' not found, using {selected_template.name}")
    
    # Check if template file exists and is readable
    if not selected_template.exists():
        logger.warning(f"Template file {selected_template} not found, using minimal built-in layout")
        readiness_report['warnings'].append("Template file not found, using built-in fallback")
        readiness_report['template_used'] = "built-in"
        return _create_minimal_presentation(), readiness_report
    
    logger.info(f"Selected template: {selected_template}")
    readiness_report['template_used'] = selected_template.name
    
    # Load presentation
    if Presentation:
        prs = Presentation(str(selected_template))
    else:
        prs = _create_blank_presentation()
    
    # Analyze layouts
    if hasattr(prs, 'slide_layouts'):
        for layout in prs.slide_layouts:
            layout_name = layout.name
            readiness_report['layouts_available'].append(layout_name)
            
            # Check for content placeholders
            content_count = _count_content_placeholders(layout)
            
            if content_count == 1:
                readiness_report['single_content_layouts'].append(layout_name)
            elif content_count == 2:
                readiness_report['dual_content_layouts'].append(layout_name)
                if config.features.single_content_only:
                    readiness_report['warnings'].append(f"Layout '{layout_name}' has dual content but single_content_only=True")
    
    # Check for required layouts
    required_layouts = ['Title Slide', 'Title and Content', 'Section Header', 'Blank']
    for required in required_layouts:
        if required not in readiness_report['layouts_available']:
            readiness_report['missing_layouts'].append(required)
    
    # Log readiness report
    logger.info(f"Template readiness: {len(readiness_report['single_content_layouts'])} single-content, "
               f"{len(readiness_report['dual_content_layouts'])} dual-content layouts")
    
    return prs, readiness_report


def _has_body_placeholder(layout: Any) -> bool:
    """Check if layout has a BODY or CONTENT placeholder."""
    try:
        pp_ph_type = None
        try:
            from pptx.enum.shapes import PP_PLACEHOLDER as _PP
            pp_ph_type = _PP
        except Exception:
            return False
            
        for ph in layout.placeholders:
            if hasattr(ph.placeholder_format, 'type'):
                valid = {pp_ph_type.BODY}
                if hasattr(pp_ph_type, "CONTENT"):
                    valid.add(pp_ph_type.CONTENT)
                if ph.placeholder_format.type in valid:
                    return True
    except Exception:
        pass
    return False


# Expanded banned layouts so insight slides never use cover-style bodies
BANNED_FOR_INSIGHT = {"cover", "title slide", "light blue cover", "section cover", "agenda", "appendix", "thank", "title only", "blank", "section"}

def _is_banned_for_insight(name: str) -> bool:
    """Check if layout name is unsuitable for insight slides."""
    return any(x in (name or "").lower() for x in BANNED_FOR_INSIGHT)


def fuzzy_map_layout(prs: Any,
                    desired_layout: str,
                    config: SynthesisConfig) -> Any:
    """
    Pick a layout with a body/content placeholder for charts.
    Avoid Agenda/Appendix/Thank You/Title Only/Blank/Cover/Section for insight slides.
    
    Args:
        prs: Presentation object
        desired_layout: Desired layout name
        config: Synthesis configuration
    
    Returns:
        Best matching layout
    """
    if not hasattr(prs, 'slide_layouts'):
        return None

    requested_lc = (desired_layout or "").lower()
    single_content_only = getattr(config.features, "single_content_only", True)

    def has_body(layout):
        return _has_body_placeholder(layout)

    def banned(name: str) -> bool:
        return _is_banned_for_insight(name)  # already lower-cases + bans cover/section/etc.

    # 1) Exact-ish match (contains)
    for layout in prs.slide_layouts:
        name_lc = layout.name.lower()
        if requested_lc and requested_lc in name_lc and has_body(layout) and not banned(layout.name):
            if single_content_only and "two content" in name_lc:
                continue
            logger.info(f"Using requested layout: {layout.name}")
            return layout

    # 2) Preferred aliases (brand-friendly)
    preferred_aliases = [
        "single column", "single content", "content only", "insight",
        "title & content", "title and content", "content with caption"
    ]
    for alias in preferred_aliases:
        for layout in prs.slide_layouts:
            name_lc = layout.name.lower()
            if alias in name_lc and has_body(layout) and not banned(layout.name):
                if single_content_only and "two content" in name_lc:
                    continue
                logger.info(f"Using preferred layout: {layout.name} (matched '{alias}')")
                return layout

    # 3) Best fallback: pick the non-banned single-content layout with the largest BODY area
    candidates = []
    for layout in prs.slide_layouts:
        if banned(layout.name) or not has_body(layout):
            continue
        # estimate body bbox area
        area = 0
        try:
            for ph in layout.placeholders:
                pf = getattr(ph, "placeholder_format", None)
                if pf and getattr(pf, "type", None):
                    # BODY/CONTENT already filtered by has_body(); use the first with width*height
                    if hasattr(ph, "width") and hasattr(ph, "height"):
                        area = max(area, int(ph.width) * int(ph.height))
        except Exception:
            pass
        candidates.append((area, layout))

    if candidates:
        candidates.sort(reverse=True, key=lambda t: t[0])
        logger.info(f"Using fallback layout with largest body: {candidates[0][1].name}")
        return candidates[0][1]

    # 4) Last resort: first non-agenda
    for layout in prs.slide_layouts:
        if "agenda" not in layout.name.lower():
            logger.warning(f"Using last-resort layout: {layout.name}")
            return layout

    logger.warning("No suitable layout found; using first layout")
    return prs.slide_layouts[0]


def add_cover_slide(prs: Any,
                   title: str,
                   subtitle: str,
                   config: SynthesisConfig,
                   as_of: Optional[str] = None) -> Any:
    """Add cover slide to presentation with optional as-of date."""
    layout = fuzzy_map_layout(prs, 'Title Slide', config)
    if not layout:
        return None
    
    slide = prs.slides.add_slide(layout)
    
    # Set title and subtitle
    if slide.shapes.title:
        slide.shapes.title.text = title
    
    # Find subtitle placeholder
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:  # Subtitle is usually index 1
            shape.text = subtitle
            break
    
    # Add date - prefer as_of if provided, otherwise use current date
    date_str = as_of or datetime.now().strftime("%B %Y")
    if as_of and not as_of.startswith("As of"):
        date_str = f"As of {as_of}"
    
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.text == "":
            shape.text = date_str
            break
    
    logger.info(f"Added cover slide with date: {date_str}")
    return slide


def add_agenda_slide(prs: Any,
                    topics: List[str],
                    config: SynthesisConfig) -> Any:
    """
    Add agenda slide with EXACT caveats.
    CRITICAL: Must include all 5 verbatim caveats.
    
    Args:
        prs: Presentation object
        topics: List of topics to include
        config: Synthesis configuration
    
    Returns:
        Agenda slide
    """
    layout = fuzzy_map_layout(prs, 'Title and Content', config)
    if not layout:
        return None
    
    slide = prs.slides.add_slide(layout)
    
    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = "Agenda"
    
    # Build agenda content
    agenda_items = []
    
    # Add topics
    for i, topic in enumerate(topics, 1):
        agenda_items.append(f"{i}. {topic.replace('_', ' ').title()}")

    if not getattr(config.features, "agenda_sections_only", False):
        # Add blank line and caveats when flag disabled
        agenda_items.append("")
        agenda_items.append("Caveats:")
        for caveat in AGENDA_CAVEATS:
            agenda_items.append(f"• {caveat}")
    
    # Find content placeholder and add text
    for shape in slide.placeholders:
        if shape.placeholder_format.idx > 0:  # Not title
            text_frame = shape.text_frame
            text_frame.clear()
            
            for item in agenda_items:
                p = text_frame.add_paragraph()
                p.text = item
                if item.startswith("•"):
                    p.level = 1
    
    _clear_default_text_artifacts(slide)
    logger.info(f"Added agenda slide with {len(topics)} topics and {len(AGENDA_CAVEATS)} caveats")
    return slide


def _find_layout_by_name(prs: Any, name: str) -> Any:
    """Find a layout by name (case-insensitive)."""
    if not hasattr(prs, 'slide_layouts'):
        return None
    name_lower = name.lower()
    for layout in prs.slide_layouts:
        if name_lower in layout.name.lower():
            return layout
    return None


def _largest_body_layout(prs: Any, banned_keywords: set = None) -> Any:
    """Find the layout with the largest body/content area, excluding banned keywords."""
    if not hasattr(prs, 'slide_layouts'):
        return None
    
    banned = banned_keywords or set()
    candidates = []
    
    for layout in prs.slide_layouts:
        # Skip if banned
        if any(k in layout.name.lower() for k in banned):
            continue
        
        # Check if has body placeholder
        if not _has_body_placeholder(layout):
            continue
            
        # Estimate body area
        area = 0
        try:
            for ph in layout.placeholders:
                pf = getattr(ph, "placeholder_format", None)
                if pf and getattr(pf, "type", None):
                    if hasattr(ph, "width") and hasattr(ph, "height"):
                        area = max(area, int(ph.width) * int(ph.height))
        except Exception:
            pass
        
        candidates.append((area, layout))
    
    if candidates:
        candidates.sort(reverse=True, key=lambda t: t[0])
        return candidates[0][1]
    
    return None


def add_section_divider(prs: Any, title: str, subtitle: str = "", config: Optional[SynthesisConfig] = None, prefer_blue_cover: bool = True) -> Any:
    """
    Add a section divider slide using a Section Header / Cover-ish layout.
    Use ONLY for topic or insight-group separators.
    
    Args:
        prs: Presentation object
        title: Section title
        subtitle: Section subtitle
        config: Synthesis configuration
        prefer_blue_cover: Whether to prefer blue cover layouts for dividers
    
    Returns:
        Section divider slide
    """
    if not prs:  # safety
        return None
    
    # Try section header / divider layouts first
    layout = None
    for cand in ("Section Header", "Section Break", "Divider", "Title Only"):
        layout = _find_layout_by_name(prs, cand)
        if layout: 
            break
    
    # Explicitly choose blue cover if allowed & requested
    if prefer_blue_cover and not layout:
        layout = _find_layout_by_name(prs, "Light blue cover")
        if layout:
            logger.info(f"Using blue cover layout for section divider")
    
    if not layout:
        # Last resort: pick largest non-banned body layout
        layout = _largest_body_layout(prs, banned_keywords={"cover"})
        logger.info(f"Using fallback divider layout: {layout.name if layout else 'None'}")
    
    if not layout:
        logger.warning("Could not find suitable layout for section divider")
        return None

    slide = prs.slides.add_slide(layout)
    _clear_default_text_artifacts(slide)

    # Set title
    target = getattr(slide.shapes, 'title', None)
    if target and getattr(target, 'text_frame', None):
        target.text_frame.clear()
        target.text_frame.paragraphs[0].text = title
    else:
        for ph in slide.placeholders:
            if getattr(ph, 'has_text_frame', False):
                try:
                    ph.text_frame.clear()
                    ph.text_frame.paragraphs[0].text = title
                except Exception:
                    pass
                break

    # Try to use a subtitle placeholder if present
    if subtitle:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx > 0:
                try:
                    ph.text_frame.clear()
                    ph.text_frame.paragraphs[0].text = subtitle
                except Exception:
                    pass
                break
    
    logger.info(f"Added section divider: {title}")
    return slide


def add_insight_slide(prs: Any,
                     insight_title: str,
                     bullets: List[str],
                     strapline: str,
                     figure_bytes: Optional[bytes],
                     footer_type: str,
                     config: SynthesisConfig,
                     chart_title: str = "",
                     notes_text: str = "") -> Any:
    """Add insight slide mapping title to insight and caption to chart title."""
    assert config.features.single_content_only, "single_content_only must be True"

    layout = fuzzy_map_layout(prs, 'Title and Content', config)
    if not layout:
        return None

    slide = prs.slides.add_slide(layout)
    _clear_default_text_artifacts(slide)

    title_ph = subtitle_ph = body_ph = content_ph = None
    for shp in slide.shapes:
        if hasattr(shp, 'placeholder_format') and PP_PLACEHOLDER:
            ph_type = shp.placeholder_format.type
            if ph_type == PP_PLACEHOLDER.TITLE:
                title_ph = shp
            elif ph_type == PP_PLACEHOLDER.SUBTITLE:
                subtitle_ph = shp
            elif ph_type == PP_PLACEHOLDER.BODY:
                body_ph = shp
            if ph_type != PP_PLACEHOLDER.TITLE and content_ph is None:
                content_ph = shp

    # Title placeholder ← insight text
    if title_ph and getattr(title_ph, 'has_text_frame', False):
        title_ph.text_frame.clear()
        title_ph.text_frame.text = insight_title
    elif slide.shapes.title:
        slide.shapes.title.text = insight_title

    # Subtitle placeholder ← strapline
    if subtitle_ph and getattr(subtitle_ph, 'has_text_frame', False) and strapline:
        subtitle_ph.text_frame.clear()
        subtitle_ph.text_frame.text = strapline

    # Caption placeholder ← chart title (if available and different from title)
    caption_target = None
    if chart_title and chart_title.strip() != (insight_title or "").strip():
        if body_ph is not None and body_ph is not content_ph and getattr(body_ph, 'has_text_frame', False):
            caption_target = body_ph
        elif subtitle_ph and not strapline and getattr(subtitle_ph, 'has_text_frame', False):
            caption_target = subtitle_ph
    if caption_target:
        caption_target.text_frame.clear()
        caption_target.text_frame.text = chart_title

    # Add footer before chart to avoid overlap
    footer_text = DECK_FOOTERS.get(footer_type, DECK_FOOTERS['default'])
    _add_footer(slide, footer_text)

    # Insert chart image with reserved footer margin
    if content_ph and figure_bytes:
        try:
            if hasattr(content_ph, 'text_frame'):
                content_ph.text_frame.clear()
        except Exception:
            pass
        png_bytes = _ensure_png_bytes(figure_bytes, dpi=config.render.figure_dpi)
        if png_bytes:
            image_stream = io.BytesIO(png_bytes)
            if hasattr(content_ph, 'insert_picture'):
                pic = content_ph.insert_picture(image_stream)
                try:
                    if hasattr(pic, 'height') and hasattr(content_ph, 'height'):
                        margin = Inches(0.35) if Inches else int(0.35 * 914400)
                        pic.height = max(0, content_ph.height - margin)
                except Exception:
                    pass
            else:
                margin = Inches(0.35) if Inches else int(0.35 * 914400)
                slide.shapes.add_picture(image_stream, content_ph.left, content_ph.top,
                                        content_ph.width, max(0, content_ph.height - margin))
            logger.info("Added chart using content placeholder")
    elif figure_bytes:
        png_bytes = _ensure_png_bytes(figure_bytes, dpi=config.render.figure_dpi)
        if png_bytes:
            image_stream = io.BytesIO(png_bytes)
            slide.shapes.add_picture(image_stream, Inches(1) if Inches else 914400,
                                    Inches(1.2) if Inches else 1097280)

    # Speaker notes
    try:
        if hasattr(slide, 'notes_slide'):
            notes = slide.notes_slide.notes_text_frame
            notes.clear()
            if notes_text:
                notes.text = notes_text
            elif bullets:
                notes.text = insight_title + "\n" + "\n".join(f"• {b}" for b in bullets)
                if strapline:
                    notes.text += f"\n\nKey insight: {strapline}"
    except Exception as e:
        logger.debug(f"Could not add speaker notes: {e}")

    # Style title
    if slide.shapes.title and hasattr(slide.shapes.title, 'text_frame'):
        try:
            for para in slide.shapes.title.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(20) if Pt else None
                    run.font.bold = True
        except Exception:
            pass

    # QA check for placeholder cleanliness
    try:
        title_count = sum(1 for shp in slide.shapes if hasattr(shp, 'placeholder_format') and PP_PLACEHOLDER and shp.placeholder_format.type == PP_PLACEHOLDER.TITLE and getattr(shp, 'text', '').strip())
        ghosts = []
        for shp in slide.shapes:
            if getattr(shp, 'has_text_frame', False) and not hasattr(shp, 'placeholder_format'):
                txt = (shp.text_frame.text or '').strip().lower()
                if txt in {"layout", "header", "subheader"}:
                    ghosts.append(txt)
        if title_count != 1 or ghosts:
            logger.warning(f"Slide text placeholders anomaly: titles={title_count}, ghosts={ghosts}")
    except Exception:
        pass

    logger.info(f"Insight header='{insight_title}' | caption='{chart_title}'")
    return slide


def add_topic_summary(prs: Any,
                     topic: str,
                     summary_text: str,
                     config: SynthesisConfig) -> Any:
    """Add topic summary slide."""
    layout = fuzzy_map_layout(prs, 'Section Header', config)
    if not layout:
        return None
    
    slide = prs.slides.add_slide(layout)
    
    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = f"{topic.replace('_', ' ').title()} Summary"
    
    # Add summary text
    for shape in slide.placeholders:
        if shape.placeholder_format.idx > 0:  # Not title
            shape.text = summary_text
            break
    
    logger.info(f"Added topic summary for {topic}")
    return slide


def add_appendix_slide(prs: Any,
                      appendix_content: Dict[str, Any],
                      config: SynthesisConfig) -> Any:
    """Add appendix slide with additional details."""
    layout = fuzzy_map_layout(prs, 'Title and Content', config)
    if not layout:
        return None
    
    slide = prs.slides.add_slide(layout)
    
    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = "Appendix"
    
    # Add content
    content_lines = []
    
    if 'data_sources' in appendix_content:
        content_lines.append("Data Sources:")
        for source in appendix_content['data_sources']:
            content_lines.append(f"• {source}")
    
    if 'methodology' in appendix_content:
        content_lines.append("")
        content_lines.append("Methodology:")
        content_lines.append(appendix_content['methodology'])
    
    if 'glossary' in appendix_content:
        content_lines.append("")
        content_lines.append("Glossary:")
        for term, definition in appendix_content['glossary'].items():
            content_lines.append(f"• {term}: {definition}")
    
    # Add to slide
    for shape in slide.placeholders:
        if shape.placeholder_format.idx > 0:  # Content placeholder
            text_frame = shape.text_frame
            text_frame.clear()
            
            for line in content_lines:
                p = text_frame.add_paragraph()
                p.text = line
                if line.startswith("•"):
                    p.level = 1
            break
    
    logger.info("Added appendix slide")
    return slide


def add_thank_you_slide(prs: Any,
                       config: SynthesisConfig,
                       client_name: Optional[str] = None) -> Any:
    """
    Add thank-you slide after appendix.
    Per dev plan requirement.
    
    Args:
        prs: Presentation object
        config: Synthesis configuration
        client_name: Optional client name for customization
    
    Returns:
        Slide object or None
    """
    # Try to find a title-only layout for clean thank-you slide
    layout = fuzzy_map_layout(prs, 'Title Slide', config)
    if not layout:
        # Fallback to title and content
        layout = fuzzy_map_layout(prs, 'Title and Content', config)
    
    if not layout:
        logger.warning("Could not find suitable layout for thank-you slide")
        return None
    
    slide = prs.slides.add_slide(layout)
    
    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = "Thank You"
    
    # Add optional subtitle or content
    try:
        # Look for subtitle placeholder
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                if client_name:
                    shape.text = f"Thank you for reviewing the {client_name} Portfolio Analysis"
                else:
                    shape.text = "Questions and feedback welcome"
                break
    except:
        # No subtitle available
        pass
    
    logger.info("Added thank-you slide")
    return slide


def export_pptx(prs: Any,
               output_path: Union[str, Path],
               config: SynthesisConfig) -> Path:
    """
    Export presentation with naming convention.
    
    Args:
        prs: Presentation object
        output_path: Output file path
        config: Synthesis configuration
    
    Returns:
        Path to exported file
    """
    output_path = Path(output_path)
    
    # Ensure .pptx extension
    if output_path.suffix != '.pptx':
        output_path = output_path.with_suffix('.pptx')
    
    # Create output directory if needed
    output_path.parent.mkdir(parents=True, exist_ok=True)
    
    # Save presentation
    if hasattr(prs, 'save'):
        prs.save(str(output_path))
        logger.info(f"Exported presentation to {output_path}")
    else:
        # python-pptx not available - raise informative error
        raise ImportError(
            "python-pptx library is required for PowerPoint export. "
            "Install with: pip install python-pptx"
        )
    
    return output_path


def export_pdf(pptx_path: Union[str, Path],
              config: SynthesisConfig,
              pdf_path: Optional[Union[str, Path]] = None) -> Optional[Path]:
    """
    Real PDF export with OS-aware fallbacks.
    
    Args:
        pptx_path: Path to PPTX file
        config: Synthesis configuration
        pdf_path: Output PDF path (auto-generate if None)
    
    Returns:
        Path to PDF or None if not enabled
    """
    if not config.features.enable_pdf_export:
        logger.info("PDF export not enabled")
        return None
    
    pptx_path = Path(pptx_path).resolve()
    pdf_path = Path(pdf_path or pptx_path.with_suffix('.pdf')).resolve()
    
    system = platform.system().lower()
    
    try:
        if system == 'windows':
            # Requires Office installed
            try:
                import comtypes.client
                powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
                powerpoint.Visible = 1
                presentation = powerpoint.Presentations.Open(str(pptx_path))
                # 32 = ppSaveAsPDF
                presentation.SaveAs(str(pdf_path), 32)
                presentation.Close()
                powerpoint.Quit()
            except ImportError:
                raise RuntimeError("comtypes not installed. Run: pip install comtypes")
                
        elif system == 'darwin':
            # Use AppleScript to drive PowerPoint if installed
            osa = shutil.which("osascript")
            if not osa:
                raise RuntimeError("osascript not found; cannot export PDF on macOS.")
            script = f'''
            tell application "Microsoft PowerPoint"
              open POSIX file "{pptx_path}"
              set theDoc to active presentation
              save as theDoc file name POSIX file "{pdf_path}" file format save as PDF
              close theDoc saving no
            end tell
            '''
            subprocess.run([osa, "-e", script], check=True)
            
        else:
            # Linux: use LibreOffice/soffice
            soffice = shutil.which("soffice") or shutil.which("libreoffice")
            if not soffice:
                raise RuntimeError("LibreOffice (soffice) not found; install libreoffice to enable PDF export.")
            outdir = str(pdf_path.parent)
            subprocess.run([soffice, "--headless", "--convert-to", "pdf", "--outdir", outdir, str(pptx_path)],
                          check=True)
            # LibreOffice names output as <name>.pdf in outdir
            if not pdf_path.exists():
                # best-effort: locate the generated pdf
                candidate = pptx_path.with_suffix('.pdf').name
                cand = Path(outdir) / candidate
                if cand.exists():
                    cand.rename(pdf_path)
        
        logger.info(f"Exported PDF: {pdf_path}")
        return pdf_path
        
    except Exception as e:
        logger.error(f"PDF export failed: {e}")
        raise


def manage_footer(slide: Any, footer_type: str) -> None:
    """
    Manage slide footer using EXACT strings.
    CRITICAL: Only use the 5 defined footers.
    
    Args:
        slide: Slide object
        footer_type: Type of footer from DECK_FOOTERS
    """
    footer_text = DECK_FOOTERS.get(footer_type, DECK_FOOTERS['default'])
    _add_footer(slide, footer_text)


# Helper functions

def _create_blank_presentation() -> Any:
    """Create blank presentation as fallback."""
    if Presentation:
        return Presentation()
    else:
        # Return stub object
        class StubPresentation:
            def __init__(self):
                self.slides = StubSlides()
                self.slide_layouts = []
        
        class StubSlides:
            def __init__(self):
                self.slides = []
            
            def add_slide(self, layout):
                slide = StubSlide()
                self.slides.append(slide)
                return slide
        
        class StubSlide:
            def __init__(self):
                self.shapes = StubShapes()
                self.placeholders = []
        
        class StubShapes:
            def __init__(self):
                self.title = StubTitle()
        
        class StubTitle:
            def __init__(self):
                self.text = ""
        
        return StubPresentation()


def _create_minimal_presentation() -> Any:
    """Create minimal presentation when no template available."""
    if Presentation:
        prs = Presentation()
        # Set standard slide dimensions (16:9)
        prs.slide_width = 10 * 914400  # 10 inches in EMUs
        prs.slide_height = 5.625 * 914400  # 5.625 inches for 16:9
        
        # Add a few basic layouts if possible
        try:
            # The presentation will have default layouts
            logger.info("Created minimal presentation with default layouts")
        except:
            logger.warning("Could not customize minimal presentation")
        
        return prs
    else:
        # Fall back to blank stub
        return _create_blank_presentation()


def _count_content_placeholders(layout: Any) -> int:
    """Count content placeholders in layout."""
    if not hasattr(layout, 'placeholders'):
        return 0

    count = 0
    # Prefer type-based detection when available
    pp_ph_type = None
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER as _PP
        pp_ph_type = _PP
    except Exception:
        pp_ph_type = None

    for ph in getattr(layout, 'placeholders', []):
        if not hasattr(ph, 'placeholder_format'):
            continue

        if pp_ph_type and hasattr(ph.placeholder_format, 'type'):
            # BODY is always present; CONTENT may not exist in some python-pptx versions
            valid = {pp_ph_type.BODY}
            if hasattr(pp_ph_type, "CONTENT"):
                valid.add(pp_ph_type.CONTENT)
            if ph.placeholder_format.type in tuple(valid):
                count += 1
        else:
            # Fallback: old heuristic on idx
            idx = getattr(ph.placeholder_format, 'idx', -1)
            if 0 < idx < 10:
                count += 1

    return count


def _add_footer(slide: Any, footer_text: str) -> None:
    """Add footer text to slide."""
    if not hasattr(slide, 'shapes'):
        return
    
    # Try to find footer placeholder
    for shape in slide.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        
        # Check if it's likely a footer (bottom of slide)
        if hasattr(shape, 'top') and hasattr(slide, 'height'):
            if shape.top > slide.height * 0.9:  # Bottom 10% of slide
                shape.text = footer_text
                return
    
    # If no footer placeholder, add as text box
    if hasattr(slide.shapes, 'add_textbox'):
        left = Inches(0.5) if Inches else 36
        top = Inches(6.5) if Inches else 468  # Near bottom
        width = Inches(9) if Inches else 648
        height = Inches(0.5) if Inches else 36
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        textbox.text_frame.text = footer_text
        
        # Format as small, gray text
        if hasattr(textbox.text_frame.paragraphs[0], 'font'):
            font = textbox.text_frame.paragraphs[0].font
            font.size = Pt(10) if Pt else 10
            if RGBColor:
                font.color.rgb = RGBColor(128, 128, 128)  # Gray
            else:
                # Fallback if RGBColor not available
                try:
                    font.color.rgb = (128, 128, 128)
                except:
                    pass  # Skip color setting if not supported